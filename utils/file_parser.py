"""File parsing and context extraction for LLM requests.

Unified file classification + parsing for common Windows/data/database formats.
"""

from __future__ import annotations

import base64
import csv
import json
import re
import sqlite3
import shutil
import zipfile
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple


class FileParser:
	"""Parse various file formats and extract content for LLM context."""

	TEXT_EXTENSIONS = {
		".txt", ".log", ".env", ".ini", ".cfg", ".conf", ".toml", ".properties", ".reg",
	}
	JSON_EXTENSIONS = {".json"}
	JSONL_EXTENSIONS = {".jsonl", ".ndjson"}
	DELIMITED_EXTENSIONS = {".csv", ".tsv", ".psv"}
	XML_EXTENSIONS = {".xml"}
	YAML_EXTENSIONS = {".yaml", ".yml"}
	MARKDOWN_EXTENSIONS = {".md", ".markdown"}
	CODE_EXTENSIONS = {".py", ".js", ".ts", ".html", ".css", ".sql", ".ps1", ".bat", ".cmd"}
	IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".ico", ".tif", ".tiff"}
	PDF_EXTENSIONS = {".pdf"}
	EMAIL_EXTENSIONS = {".eml", ".msg"}
	RICH_TEXT_EXTENSIONS = {".rtf"}
	OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
	ARCHIVE_EXTENSIONS = {".zip"}
	WINDOWS_DIAGNOSTIC_EXTENSIONS = {".wer", ".dmp", ".mdmp", ".lnk"}
	SQLITE_EXTENSIONS = {
		".db", ".sqlite", ".sqlite3", ".db3", ".s3db", ".sl3", ".dbs",
	}
	SQLITE_SIDECAR_ENDINGS = {
		".db-wal", ".db-shm", ".sqlite-wal", ".sqlite-shm", ".sqlite3-wal", ".sqlite3-shm",
	}
	WINDOWS_BINARY_DATA_EXTENSIONS = {
		".dat", ".edb", ".evtx", ".etl", ".blf", ".jrs", ".chk", ".cab", ".msi",
	}
	KNOWN_SQLITE_FILENAMES = {
		"history", "cookies", "login data", "web data", "favicons", "network action predictor",
	}

	SUPPORTED_EXTENSIONS = (
		TEXT_EXTENSIONS
		| JSON_EXTENSIONS
		| JSONL_EXTENSIONS
		| DELIMITED_EXTENSIONS
		| XML_EXTENSIONS
		| YAML_EXTENSIONS
		| MARKDOWN_EXTENSIONS
		| CODE_EXTENSIONS
		| IMAGE_EXTENSIONS
		| PDF_EXTENSIONS
		| EMAIL_EXTENSIONS
		| RICH_TEXT_EXTENSIONS
		| OFFICE_EXTENSIONS
		| ARCHIVE_EXTENSIONS
		| WINDOWS_DIAGNOSTIC_EXTENSIONS
		| SQLITE_EXTENSIONS
		| WINDOWS_BINARY_DATA_EXTENSIONS
	)

	MAX_FILE_SIZE_MB = 10
	MAX_FILE_SIZE_MB_DATABASE = 500  # Databases can be larger; we only read what we need
	MAX_CONTENT_LENGTH = 50000  # chars for text, bytes for binary
	MAGIC_BYTES_READ = 64

	@staticmethod
	def _read_magic(file_path: Path, count: int = MAGIC_BYTES_READ) -> bytes:
		try:
			with open(file_path, "rb") as f:
				return f.read(count)
		except Exception:
			return b""

	@classmethod
	def _is_sqlite_magic(cls, file_path: Path) -> bool:
		magic = cls._read_magic(file_path, 16)
		return magic.startswith(b"SQLite format 3")

	@classmethod
	def _is_sqlite_sidecar(cls, file_path: Path) -> bool:
		name = file_path.name.lower()
		return any(name.endswith(ending) for ending in cls.SQLITE_SIDECAR_ENDINGS)

	@classmethod
	def _looks_like_json(cls, file_path: Path) -> bool:
		magic = cls._read_magic(file_path, 64).lstrip()
		return magic.startswith(b"{") or magic.startswith(b"[")

	@classmethod
	def _looks_like_xml(cls, file_path: Path) -> bool:
		magic = cls._read_magic(file_path, 64).lstrip()
		return magic.startswith(b"<?xml") or magic.startswith(b"<")

	@classmethod
	def _is_known_sqlite_filename(cls, file_path: Path) -> bool:
		return file_path.name.strip().lower() in cls.KNOWN_SQLITE_FILENAMES

	@classmethod
	def _looks_binary_blob(cls, file_path: Path) -> bool:
		sample = cls._read_magic(file_path, 2048)
		if not sample:
			return False
		if b"\x00" in sample:
			return True
		# If most bytes are non-printable, treat as binary.
		printable = set(range(32, 127)) | {9, 10, 13}
		non_printable = sum(1 for b in sample if b not in printable)
		return (non_printable / max(len(sample), 1)) > 0.30

	@classmethod
	def _detect_file_type(cls, file_path: Path) -> str:
		suffix = file_path.suffix.lower()
		name_lower = file_path.name.lower()

		if cls._is_sqlite_sidecar(file_path):
			return "sqlite_sidecar"
		if suffix in cls.SQLITE_EXTENSIONS:
			return "sqlite"
		if suffix in cls.JSON_EXTENSIONS:
			return "json"
		if suffix in cls.JSONL_EXTENSIONS:
			return "jsonl"
		if suffix in cls.DELIMITED_EXTENSIONS:
			if suffix == ".tsv":
				return "tsv"
			if suffix == ".psv":
				return "psv"
			return "csv"
		if suffix in cls.XML_EXTENSIONS:
			return "xml"
		if suffix in cls.YAML_EXTENSIONS:
			return "yaml"
		if suffix in cls.MARKDOWN_EXTENSIONS:
			return "markdown"
		if suffix in cls.CODE_EXTENSIONS:
			return "code"
		if suffix in cls.TEXT_EXTENSIONS:
			return "text"
		if suffix in cls.IMAGE_EXTENSIONS:
			return "image"
		if suffix in cls.PDF_EXTENSIONS:
			return "pdf"
		if suffix == ".eml":
			return "eml"
		if suffix == ".msg":
			return "msg"
		if suffix in cls.RICH_TEXT_EXTENSIONS:
			return "rtf"
		if suffix == ".docx":
			return "office_docx"
		if suffix == ".xlsx":
			return "office_xlsx"
		if suffix == ".pptx":
			return "office_pptx"
		if suffix in cls.ARCHIVE_EXTENSIONS:
			return "zip"
		if suffix == ".wer":
			return "wer"
		if suffix == ".lnk":
			return "lnk"
		if suffix in {".dmp", ".mdmp"}:
			return "memory_dump"
		if suffix in cls.WINDOWS_BINARY_DATA_EXTENSIONS:
			return "binary_data"

		if file_path.exists() and file_path.is_file():
			if cls._is_sqlite_magic(file_path):
				return "sqlite"
			if cls._looks_like_json(file_path):
				return "json"
			if cls._looks_like_xml(file_path):
				return "xml"
			if cls._is_known_sqlite_filename(file_path) and cls._is_sqlite_magic(file_path):
				return "sqlite"
			if name_lower.endswith((".wal", ".shm")):
				return "sqlite_sidecar"
			if cls._looks_binary_blob(file_path):
				return "unknown"

		return "text"

	@staticmethod
	def is_supported(file_path: Path) -> bool:
		"""Check if file type is supported by extension, filename patterns, or magic bytes."""
		if not file_path.exists() or not file_path.is_file():
			return False
		return FileParser._detect_file_type(file_path) != "unknown"

	@staticmethod
	def get_file_size_mb(file_path: Path) -> float:
		"""Get file size in MB."""
		return file_path.stat().st_size / (1024 * 1024)

	@classmethod
	def parse_file(cls, file_path: Path) -> Tuple[bool, str]:
		"""Parse file and return (success, content_or_error_message)."""
		if not file_path.exists():
			return False, f"File not found: {file_path}"

		if not file_path.is_file():
			return False, f"Not a file: {file_path}"

		if not cls.is_supported(file_path):
			return False, f"Unsupported file type: {file_path.suffix}"

		size_mb = cls.get_file_size_mb(file_path)
		file_type = cls._detect_file_type(file_path)

		# Use larger limit for databases
		max_size = cls.MAX_FILE_SIZE_MB_DATABASE if file_type in {"sqlite", "sqlite_sidecar"} else cls.MAX_FILE_SIZE_MB
		if size_mb > max_size:
			return False, f"File too large: {size_mb:.1f}MB (max {max_size}MB)"

		try:
			if file_type == "json":
				return cls._parse_json(file_path)
			elif file_type == "csv":
				return cls._parse_csv(file_path)
			elif file_type == "tsv":
				return cls._parse_delimited(file_path, delimiter="\t", label="TSV")
			elif file_type == "psv":
				return cls._parse_delimited(file_path, delimiter="|", label="PSV")
			elif file_type == "jsonl":
				return cls._parse_json_lines(file_path)
			elif file_type == "xml":
				return cls._parse_xml(file_path)
			elif file_type == "yaml":
				return cls._parse_yaml(file_path)
			elif file_type == "markdown":
				return cls._parse_markdown(file_path)
			elif file_type == "code":
				return cls._parse_code(file_path)
			elif file_type == "text":
				return cls._parse_text(file_path)
			elif file_type == "image":
				return cls._parse_image(file_path)
			elif file_type == "pdf":
				return cls._parse_pdf(file_path)
			elif file_type == "eml":
				return cls._parse_eml(file_path)
			elif file_type == "msg":
				return cls._parse_msg(file_path)
			elif file_type == "rtf":
				return cls._parse_rtf(file_path)
			elif file_type == "office_docx":
				return cls._parse_docx(file_path)
			elif file_type == "office_xlsx":
				return cls._parse_xlsx(file_path)
			elif file_type == "office_pptx":
				return cls._parse_pptx(file_path)
			elif file_type == "zip":
				return cls._parse_zip_manifest(file_path)
			elif file_type == "wer":
				return cls._parse_wer(file_path)
			elif file_type == "lnk":
				return cls._parse_lnk(file_path)
			elif file_type == "memory_dump":
				return cls._parse_memory_dump(file_path)
			elif file_type == "sqlite":
				return cls._parse_sqlite(file_path)
			elif file_type == "sqlite_sidecar":
				return cls._parse_sqlite_sidecar(file_path)
			elif file_type == "binary_data":
				return cls._parse_binary_metadata(file_path)
			else:
				return cls._parse_text(file_path)
		except Exception as e:
			return False, f"Parse error ({type(e).__name__}): {e}"

	@staticmethod
	def _parse_json(file_path: Path) -> Tuple[bool, str]:
		"""Parse JSON file with pretty formatting."""
		with open(file_path, "r", encoding="utf-8") as f:
			data = json.load(f)
		formatted = json.dumps(data, indent=2)
		if len(formatted) > FileParser.MAX_CONTENT_LENGTH:
			formatted = formatted[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
		return True, f"```json\n{formatted}\n```"

	@staticmethod
	def _parse_csv(file_path: Path) -> Tuple[bool, str]:
		"""Parse CSV file with column headers."""
		with open(file_path, "r", encoding="utf-8", errors="replace") as f:
			reader = csv.DictReader(f)
			rows = list(reader)

		if not rows:
			return True, "CSV file is empty"

		header = list(rows[0].keys()) if rows else []
		content = f"CSV Columns: {', '.join(header)}\n\nRows:\n"
		content += json.dumps(rows[: 100], indent=2)  # First 100 rows

		if len(rows) > 100:
			content += f"\n...[{len(rows) - 100} more rows]"

		if len(content) > FileParser.MAX_CONTENT_LENGTH:
			content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"

		return True, f"```\n{content}\n```"

	@staticmethod
	def _parse_delimited(file_path: Path, delimiter: str, label: str) -> Tuple[bool, str]:
		"""Parse delimiter-separated text files (TSV/PSV)."""
		with open(file_path, "r", encoding="utf-8", errors="replace") as f:
			reader = csv.reader(f, delimiter=delimiter)
			rows = list(reader)

		if not rows:
			return True, f"{label} file is empty"

		header = rows[0]
		data_rows = rows[1:101]
		content = f"{label} Columns: {', '.join(header)}\n\nRows:\n"
		content += json.dumps(data_rows, indent=2)

		if len(rows) > 101:
			content += f"\n...[{len(rows) - 101} more rows]"

		if len(content) > FileParser.MAX_CONTENT_LENGTH:
			content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"

		return True, f"```\n{content}\n```"

	@staticmethod
	def _parse_json_lines(file_path: Path) -> Tuple[bool, str]:
		"""Parse JSON Lines / NDJSON files."""
		rows: List[Any] = []
		errors = 0
		with open(file_path, "r", encoding="utf-8", errors="replace") as f:
			for idx, line in enumerate(f):
				if idx >= 200:
					break
				line = line.strip()
				if not line:
					continue
				try:
					rows.append(json.loads(line))
				except Exception:
					errors += 1

		if not rows and errors:
			return False, "NDJSON parse failed: no valid JSON lines"

		content = f"Parsed JSONL records: {len(rows)}"
		if errors:
			content += f" (invalid lines skipped: {errors})"
		content += "\n\n"
		content += json.dumps(rows[:100], indent=2)

		if len(content) > FileParser.MAX_CONTENT_LENGTH:
			content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"

		return True, f"```json\n{content}\n```"

	@staticmethod
	def _parse_xml(file_path: Path) -> Tuple[bool, str]:
		"""Parse XML file (basic pretty-print)."""
		content = file_path.read_text(encoding="utf-8", errors="replace")
		try:
			import xml.dom.minidom as minidom
			dom = minidom.parseString(content)
			pretty = dom.toprettyxml(indent="  ")
			# Remove XML declaration for brevity
			if pretty.startswith("<?xml"):
				pretty = "\n".join(pretty.split("\n")[1:])
		except Exception:
			pretty = content

		if len(pretty) > FileParser.MAX_CONTENT_LENGTH:
			pretty = pretty[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"

		return True, f"```xml\n{pretty}\n```"

	@staticmethod
	def _parse_yaml(file_path: Path) -> Tuple[bool, str]:
		"""Parse YAML file (falls back to raw if yaml not installed)."""
		content = file_path.read_text(encoding="utf-8", errors="replace")
		try:
			import yaml
			data = yaml.safe_load(content)
			formatted = json.dumps(data, indent=2)
		except ImportError:
			formatted = content
		except Exception:
			formatted = content

		if len(formatted) > FileParser.MAX_CONTENT_LENGTH:
			formatted = formatted[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"

		return True, f"```yaml\n{formatted}\n```"

	@staticmethod
	def _parse_markdown(file_path: Path) -> Tuple[bool, str]:
		"""Parse Markdown file."""
		content = file_path.read_text(encoding="utf-8", errors="replace")
		if len(content) > FileParser.MAX_CONTENT_LENGTH:
			content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
		return True, f"```markdown\n{content}\n```"

	@staticmethod
	def _parse_code(file_path: Path) -> Tuple[bool, str]:
		"""Parse code file with syntax highlighting hint."""
		content = file_path.read_text(encoding="utf-8", errors="replace")
		lang_map = {
			".py": "python",
			".js": "javascript",
			".ts": "typescript",
			".html": "html",
			".css": "css",
			".sql": "sql",
			".ps1": "powershell",
			".bat": "bat",
			".cmd": "bat",
		}
		lang = lang_map.get(file_path.suffix.lower(), "txt")

		if len(content) > FileParser.MAX_CONTENT_LENGTH:
			content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"

		return True, f"```{lang}\n{content}\n```"

	@staticmethod
	def _parse_text(file_path: Path) -> Tuple[bool, str]:
		"""Parse plain text file."""
		content = file_path.read_text(encoding="utf-8", errors="replace")
		if len(content) > FileParser.MAX_CONTENT_LENGTH:
			content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
		return True, f"```\n{content}\n```"

	@staticmethod
	def _parse_sqlite_sidecar(file_path: Path) -> Tuple[bool, str]:
		"""Describe SQLite sidecar files (WAL/SHM) without deep parsing."""
		size_bytes = file_path.stat().st_size
		content = (
			f"SQLite sidecar file: {file_path.name}\n"
			f"Path: {file_path}\n"
			f"Size: {size_bytes} bytes\n"
			"Type: Write-ahead log/shared memory companion for a SQLite DB\n"
			"Note: Parse the base .db/.sqlite file for queryable records."
		)
		return True, f"```\n{content}\n```"

	@staticmethod
	def _parse_binary_metadata(file_path: Path) -> Tuple[bool, str]:
		"""Return metadata for known Windows binary data containers."""
		magic = FileParser._read_magic(file_path, 32)
		size_bytes = file_path.stat().st_size
		hex_magic = magic.hex(" ") if magic else "(unavailable)"
		content = (
			f"Binary data file: {file_path.name}\n"
			f"Path: {file_path}\n"
			f"Extension: {file_path.suffix.lower() or '(none)'}\n"
			f"Size: {size_bytes} bytes\n"
			f"Magic bytes (first {len(magic)}): {hex_magic}\n"
			"Note: This parser currently provides metadata only for this binary container type."
		)
		return True, f"```\n{content}\n```"

	@staticmethod
	def _parse_image(file_path: Path) -> Tuple[bool, str]:
		"""Parse image file (metadata + base64 preview)."""
		import mimetypes

		mime_type, _ = mimetypes.guess_type(str(file_path))
		if not mime_type:
			mime_type = "application/octet-stream"

		size_kb = file_path.stat().st_size / 1024
		content = f"Image file: {file_path.name}\nMIME type: {mime_type}\nSize: {size_kb:.1f}KB\n"

		# Include base64 for small images
		if size_kb < 100:  # < 100KB
			with open(file_path, "rb") as f:
				b64 = base64.b64encode(f.read()).decode("ascii")
			content += f"\n[Base64 (first 500 chars)]\n{b64[:500]}...\n"
			content += f"[Use this base64 in vision requests if needed]"
		else:
			content += f"[File too large for inline embedding; use external path]"

		return True, content

	@staticmethod
	def _parse_pdf(file_path: Path) -> Tuple[bool, str]:
		"""Parse PDF file (text extraction)."""
		try:
			import PyPDF2
			with open(file_path, "rb") as f:
				reader = PyPDF2.PdfReader(f)
				text = ""
				for page_num, page in enumerate(reader.pages[: 20]):  # First 20 pages
					text += f"\n--- Page {page_num + 1} ---\n"
					text += page.extract_text()

			if len(text) > FileParser.MAX_CONTENT_LENGTH:
				text = text[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"

			return True, f"```\n{text}\n```"
		except ImportError:
			return False, "PyPDF2 not installed. Install with: pip install PyPDF2"
		except Exception as e:
			return False, f"PDF parsing failed: {e}"

	@staticmethod
	def _parse_eml(file_path: Path) -> Tuple[bool, str]:
		"""Parse RFC email (.eml)."""
		with open(file_path, "rb") as f:
			message = BytesParser(policy=policy.default).parse(f)

		subject = str(message.get("subject", ""))
		from_addr = str(message.get("from", ""))
		to_addr = str(message.get("to", ""))
		date = str(message.get("date", ""))

		plain_parts: List[str] = []
		html_parts: List[str] = []
		attachments: List[str] = []

		for part in message.walk():
			content_disposition = str(part.get_content_disposition() or "")
			filename = part.get_filename()
			if filename:
				attachments.append(filename)
				continue

			content_type = part.get_content_type()
			if content_disposition == "attachment":
				continue
			if content_type == "text/plain":
				plain_parts.append(part.get_content())
			elif content_type == "text/html":
				html_parts.append(part.get_content())

		body = "\n\n".join([p for p in plain_parts if p]).strip()
		if not body and html_parts:
			body = "[HTML body present]"

		content = (
			f"EML Message: {file_path.name}\n"
			f"Subject: {subject}\n"
			f"From: {from_addr}\n"
			f"To: {to_addr}\n"
			f"Date: {date}\n"
		)
		if attachments:
			content += f"Attachments ({len(attachments)}): {', '.join(attachments[:20])}\n"
		content += "\nBody:\n"
		content += body[: FileParser.MAX_CONTENT_LENGTH]

		if len(content) > FileParser.MAX_CONTENT_LENGTH:
			content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"

		return True, f"```\n{content}\n```"

	@staticmethod
	def _parse_msg(file_path: Path) -> Tuple[bool, str]:
		"""Parse Outlook .msg using optional extract-msg dependency."""
		try:
			import extract_msg
		except ImportError:
			size_bytes = file_path.stat().st_size
			fallback = (
				f"MSG Message: {file_path.name}\n"
				f"Size: {size_bytes} bytes\n"
				"Detailed MSG parsing requires optional dependency: pip install extract-msg"
			)
			return True, f"```\n{fallback}\n```"

		try:
			msg = extract_msg.Message(str(file_path))
			msg_subject = msg.subject or ""
			msg_sender = msg.sender or ""
			msg_to = msg.to or ""
			msg_date = str(getattr(msg, "date", "") or "")
			msg_body = (msg.body or "")[: FileParser.MAX_CONTENT_LENGTH]
			attachments = [att.longFilename or att.filename for att in msg.attachments]

			content = (
				f"MSG Message: {file_path.name}\n"
				f"Subject: {msg_subject}\n"
				f"From: {msg_sender}\n"
				f"To: {msg_to}\n"
				f"Date: {msg_date}\n"
			)
			if attachments:
				content += f"Attachments ({len(attachments)}): {', '.join([a for a in attachments if a][:20])}\n"
			content += "\nBody:\n" + msg_body
			return True, f"```\n{content}\n```"
		except Exception as e:
			return False, f"MSG parsing failed: {e}"

	@staticmethod
	def _parse_rtf(file_path: Path) -> Tuple[bool, str]:
		"""Parse rich text format using optional striprtf dependency."""
		try:
			from striprtf.striprtf import rtf_to_text
		except ImportError:
			raw = file_path.read_text(encoding="utf-8", errors="replace")
			fallback = re.sub(r"\\[a-zA-Z]+[0-9]* ?", "", raw)
			fallback = fallback.replace("{", "").replace("}", "")
			if len(fallback) > FileParser.MAX_CONTENT_LENGTH:
				fallback = fallback[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
			return True, f"```\n{fallback}\n```"

		try:
			raw = file_path.read_text(encoding="utf-8", errors="replace")
			text = rtf_to_text(raw)
			if len(text) > FileParser.MAX_CONTENT_LENGTH:
				text = text[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
			return True, f"```\n{text}\n```"
		except Exception as e:
			return False, f"RTF parsing failed: {e}"

	@staticmethod
	def _parse_docx(file_path: Path) -> Tuple[bool, str]:
		"""Parse DOCX using optional python-docx dependency."""
		try:
			from docx import Document
		except ImportError:
			fallback = (
				f"DOCX document: {file_path.name}\n"
				"Detailed DOCX parsing requires optional dependency: pip install python-docx"
			)
			return True, f"```\n{fallback}\n```"

		try:
			doc = Document(str(file_path))
			paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
			content = "\n".join(paragraphs)
			if len(content) > FileParser.MAX_CONTENT_LENGTH:
				content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
			return True, f"```\n{content}\n```"
		except Exception as e:
			return False, f"DOCX parsing failed: {e}"

	@staticmethod
	def _parse_xlsx(file_path: Path) -> Tuple[bool, str]:
		"""Parse XLSX using optional openpyxl dependency."""
		try:
			import openpyxl
		except ImportError:
			fallback = (
				f"XLSX workbook: {file_path.name}\n"
				"Detailed XLSX parsing requires optional dependency: pip install openpyxl"
			)
			return True, f"```\n{fallback}\n```"

		try:
			wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
			lines: List[str] = [f"Workbook: {file_path.name}", f"Sheets: {', '.join(wb.sheetnames)}", ""]
			for sheet_name in wb.sheetnames[:5]:
				ws = wb[sheet_name]
				lines.append(f"--- Sheet: {sheet_name} ---")
				for idx, row in enumerate(ws.iter_rows(min_row=1, max_row=25, values_only=True), 1):
					vals = ["" if v is None else str(v) for v in row]
					lines.append(f"{idx:>3}: " + " | ".join(vals[:40]))
				lines.append("")
			content = "\n".join(lines)
			if len(content) > FileParser.MAX_CONTENT_LENGTH:
				content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
			return True, f"```\n{content}\n```"
		except Exception as e:
			return False, f"XLSX parsing failed: {e}"

	@staticmethod
	def _parse_pptx(file_path: Path) -> Tuple[bool, str]:
		"""Parse PPTX using optional python-pptx dependency."""
		try:
			from pptx import Presentation
		except ImportError:
			fallback = (
				f"PPTX presentation: {file_path.name}\n"
				"Detailed PPTX parsing requires optional dependency: pip install python-pptx"
			)
			return True, f"```\n{fallback}\n```"

		try:
			prs = Presentation(str(file_path))
			lines: List[str] = [f"Presentation: {file_path.name}", f"Slides: {len(prs.slides)}", ""]
			for slide_idx, slide in enumerate(prs.slides[:50], 1):
				lines.append(f"--- Slide {slide_idx} ---")
				for shape in slide.shapes:
					if hasattr(shape, "text") and shape.text:
						text = shape.text.strip()
						if text:
							lines.append(text)
			content = "\n".join(lines)
			if len(content) > FileParser.MAX_CONTENT_LENGTH:
				content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
			return True, f"```\n{content}\n```"
		except Exception as e:
			return False, f"PPTX parsing failed: {e}"

	@staticmethod
	def _parse_zip_manifest(file_path: Path) -> Tuple[bool, str]:
		"""Parse ZIP manifest and preview parsable entries."""
		try:
			with zipfile.ZipFile(file_path, "r") as zf:
				names = zf.namelist()
				entries = [n for n in names if not n.endswith("/")]
				lines: List[str] = [
					f"ZIP Archive: {file_path.name}",
					f"Entries: {len(entries)}",
					"",
					"Manifest (first 200):",
				]
				for entry in entries[:200]:
					try:
						info = zf.getinfo(entry)
						lines.append(f"- {entry} ({info.file_size} bytes)")
					except Exception:
						lines.append(f"- {entry}")

				preview_candidates = [
					e for e in entries
					if Path(e).suffix.lower() in {".txt", ".log", ".json", ".csv", ".xml", ".yaml", ".yml", ".md"}
				][:5]
				if preview_candidates:
					lines.append("")
					lines.append("Inline previews:")
					for entry in preview_candidates:
						try:
							raw = zf.read(entry)[:2000]
							text = raw.decode("utf-8", errors="replace")
							lines.append(f"-- {entry} --")
							lines.append(text)
						except Exception:
							continue

				content = "\n".join(lines)
				if len(content) > FileParser.MAX_CONTENT_LENGTH:
					content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
				return True, f"```\n{content}\n```"
		except zipfile.BadZipFile:
			return False, "Invalid ZIP archive"
		except Exception as e:
			return False, f"ZIP parsing failed: {e}"

	@staticmethod
	def _parse_wer(file_path: Path) -> Tuple[bool, str]:
		"""Parse Windows Error Reporting (.wer) key/value content."""
		content = ""
		for enc in ("utf-16", "utf-8", "latin-1"):
			try:
				candidate = file_path.read_text(encoding=enc, errors="replace")
				if candidate:
					content = candidate
					if "=" in candidate:
						break
			except Exception:
				continue
		if not content:
			return False, "WER parsing failed: could not decode file"

		lines = [line.strip() for line in content.splitlines() if line.strip()]
		parsed: Dict[str, str] = {}
		for line in lines:
			if "=" in line:
				k, v = line.split("=", 1)
				parsed[k.strip()] = v.strip()

		if not parsed:
			if len(content) > FileParser.MAX_CONTENT_LENGTH:
				content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
			return True, f"```\n{content}\n```"

		formatted = json.dumps(parsed, indent=2)
		if len(formatted) > FileParser.MAX_CONTENT_LENGTH:
			formatted = formatted[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"
		return True, f"```json\n{formatted}\n```"

	@staticmethod
	def _parse_lnk(file_path: Path) -> Tuple[bool, str]:
		"""Parse Windows shortcut (.lnk) with optional pylnk3 dependency."""
		try:
			import pylnk3
		except ImportError:
			size_bytes = file_path.stat().st_size
			fallback = (
				f"Windows Shortcut: {file_path.name}\n"
				f"Size: {size_bytes} bytes\n"
				"Install optional parser for detailed fields: pip install pylnk3"
			)
			return True, f"```\n{fallback}\n```"

		try:
			lnk = pylnk3.parse(str(file_path))
			content = {
				"name": file_path.name,
				"path": str(file_path),
				"target": getattr(lnk, "path", None),
				"arguments": getattr(lnk, "arguments", None),
				"working_dir": getattr(lnk, "work_dir", None),
				"description": getattr(lnk, "description", None),
				"icon": getattr(lnk, "icon", None),
			}
			return True, f"```json\n{json.dumps(content, indent=2)}\n```"
		except Exception as e:
			return False, f"LNK parsing failed: {e}"

	@staticmethod
	def _parse_memory_dump(file_path: Path) -> Tuple[bool, str]:
		"""Parse memory dump metadata only."""
		size_bytes = file_path.stat().st_size
		magic = FileParser._read_magic(file_path, 32)
		hex_magic = magic.hex(" ") if magic else "(unavailable)"
		content = (
			f"Memory dump file: {file_path.name}\n"
			f"Path: {file_path}\n"
			f"Size: {size_bytes} bytes\n"
			f"Magic bytes: {hex_magic}\n"
			"Note: Full dump analysis requires external debuggers (WinDbg/cdb)."
		)
		return True, f"```\n{content}\n```"

	@staticmethod
	def _parse_sqlite(file_path: Path) -> Tuple[bool, str]:
		"""Parse SQLite database file (schema inspection + sample data)."""
		# Make a temporary copy since Chrome History database may be locked
		temp_copy = None
		try:
			# If this is the Chrome History file, make a copy to avoid locking issues
			if "Chrome" in str(file_path) and "History" in file_path.name:
				import tempfile
				temp_dir = tempfile.gettempdir()
				temp_copy = Path(temp_dir) / f"chrome_history_temp_{hash(file_path) % 100000}.db"
				shutil.copy2(file_path, temp_copy)
				db_path = temp_copy
			else:
				db_path = file_path

			conn = sqlite3.connect(str(db_path))
			cursor = conn.cursor()

			# Get all table names
			cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
			tables = [row[0] for row in cursor.fetchall()]

			if not tables:
				conn.close()
				return True, "SQLite database is empty (no tables)"

			output = f"SQLite Database: {file_path.name}\nTables: {', '.join(tables)}\n\n"

			# Special handling for Chrome History database
			if "Chrome" in str(file_path) and "History" in file_path.name:
				output += "=== CHROME HISTORY ANALYSIS ===\n\n"

				# Chrome stores timestamps as microseconds since Windows epoch (1601-01-01)
				# Conversion: Google Chrome uses Windows FILETIME format
				# May 17, 2026 00:00:00 UTC = 133629312000000000 (100-nanosecond intervals from 1601-01-01)
				# May 18, 2026 23:59:59 UTC = 133630463999999999
				# In microseconds (divide by 10):
				may_17_start = 13362931200000000  # May 17, 2026 00:00:00 UTC
				may_18_end = 13363046399999999     # May 18, 2026 23:59:59 UTC

				try:
					cursor.execute("SELECT COUNT(*) FROM urls")
					url_count = cursor.fetchone()[0]
					output += f"Total URLs in history: {url_count}\n\n"

					# Get URLs with visit count
					cursor.execute("""
						SELECT url, title, visit_count 
						FROM urls 
						ORDER BY visit_count DESC 
						LIMIT 20
					""")
					top_urls = cursor.fetchall()
					output += "Top 20 Most Visited URLs (all time):\n"
					for idx, (url, title, count) in enumerate(top_urls, 1):
						title_str = (title[:50] + "...") if title and len(title) > 50 else (title or "N/A")
						output += f"  {idx}. [{count} visits] {url[:80]}\n"
						if title:
							output += f"     Title: {title_str}\n"
				except Exception as e:
					output += f"Could not analyze URLs table: {e}\n"

				# Get visits with timestamps for the date range
				try:
					# Query visits in the May 17-18 date range
					cursor.execute(f"""
						SELECT u.url, v.visit_time, v.transition
						FROM visits v
						JOIN urls u ON v.url = u.id
						WHERE v.visit_time BETWEEN {may_17_start} AND {may_18_end}
						ORDER BY v.visit_time DESC
						LIMIT 100
					""")
					dated_visits = cursor.fetchall()
					
					if dated_visits:
						output += f"\n=== VISITS ON MAY 17-18, 2026 ===\n"
						output += f"Total visits during this period: {len(dated_visits)} (showing top 100)\n\n"
						
						# Group by domain
						domains = {}
						for url, visit_time, transition in dated_visits:
							try:
								from urllib.parse import urlparse
								domain = urlparse(url).netloc or url[:50]
								if domain not in domains:
									domains[domain] = []
								domains[domain].append((url, visit_time))
							except:
								domain = url[:50]
								if domain not in domains:
									domains[domain] = []
								domains[domain].append((url, visit_time))
						
						output += f"Unique domains visited: {len(domains)}\n\n"
						output += "Top domains by visit frequency (May 17-18):\n"
						sorted_domains = sorted(domains.items(), key=lambda x: len(x[1]), reverse=True)
						for idx, (domain, visits) in enumerate(sorted_domains[:15], 1):
							output += f"  {idx}. {domain} ({len(visits)} visits)\n"
					else:
						output += f"\nNo visits recorded for May 17-18, 2026\n"
				except Exception as e:
					output += f"Could not analyze visits by date: {e}\n"

				# Overall visits count
				try:
					cursor.execute("SELECT COUNT(*) FROM visits")
					visit_count = cursor.fetchone()[0]
					output += f"\nTotal visits in history (all time): {visit_count}\n"
				except Exception as e:
					output += f"Could not count visits: {e}\n"

			else:
				# Generic database analysis for non-Chrome databases
				for table in tables[:5]:  # First 5 tables
					try:
						cursor.execute(f"SELECT COUNT(*) FROM [{table}]")
						row_count = cursor.fetchone()[0]
						output += f"\nTable: {table} ({row_count} rows)\n"

						# Get column info
						cursor.execute(f"PRAGMA table_info([{table}])")
						columns = cursor.fetchall()
						output += f"  Columns: {', '.join([col[1] for col in columns])}\n"

						# Sample data
						cursor.execute(f"SELECT * FROM [{table}] LIMIT 5")
						sample = cursor.fetchall()
						if sample:
							output += f"  Sample data (first 5 rows):\n"
							for row in sample:
								output += f"    {row}\n"
					except Exception as e:
						output += f"  Error reading table {table}: {e}\n"

				if len(tables) > 5:
					output += f"\n... and {len(tables) - 5} more tables"

			conn.close()

			# Cleanup temp copy
			if temp_copy and temp_copy.exists():
				try:
					temp_copy.unlink()
				except:
					pass

			if len(output) > FileParser.MAX_CONTENT_LENGTH:
				output = output[: FileParser.MAX_CONTENT_LENGTH] + "\n...[truncated]"

			return True, f"```\n{output}\n```"

		except sqlite3.DatabaseError:
			return False, "Not a valid SQLite database file"
		except Exception as e:
			# Cleanup temp copy on error
			if temp_copy and temp_copy.exists():
				try:
					temp_copy.unlink()
				except:
					pass
			return False, f"SQLite parsing failed: {e}"


class FileContextBuilder:
	"""Build prompt context from multiple files."""

	@staticmethod
	def build_context(file_paths: List[Path]) -> Tuple[List[Tuple[Path, str]], List[str]]:
		"""Parse files and return (parsed_files, error_messages)."""
		parsed = []
		errors = []

		for file_path in file_paths:
			success, content = FileParser.parse_file(file_path)
			if success:
				parsed.append((file_path, content))
			else:
				errors.append(f"{file_path.name}: {content}")

		return parsed, errors

	@staticmethod
	def format_context_block(file_path: Path, content: str) -> str:
		"""Format a single file as context block."""
		return f"## File: {file_path.name}\n\n{content}"

	@staticmethod
	def combine_contexts(parsed_files: List[Tuple[Path, str]]) -> str:
		"""Combine multiple parsed files into single context string."""
		if not parsed_files:
			return ""

		blocks = [
			FileContextBuilder.format_context_block(file_path, content)
			for file_path, content in parsed_files
		]

		combined = "\n\n---\n\n".join(blocks)
		if len(combined) > 100000:  # 100K char limit for full context
			combined = combined[:100000] + "\n\n...[context truncated]"

		return combined
